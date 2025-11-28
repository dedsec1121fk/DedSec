#!/usr/bin/env python

import sys
import os
import subprocess
import importlib.util
import time
import traceback
import zipfile
import tarfile
import csv
import json
import gzip
import shutil
from contextlib import redirect_stderr, redirect_stdout

# --- 1. ΡΥΘΜΙΣΕΙΣ & ΔΙΑΜΟΡΦΩΣΗ ---

# (14) Βιβλιοθήκες Python για αυτόματη εγκατάσταση
REQUIRED_MODULES = {
    "Pillow": "Pillow",         # Εικόνες
    "reportlab": "reportlab",   # Δημιουργία PDF
    "docx": "python-docx",    # Έγγραφα Word
    "odf": "odfpy",           # Έγγραφα OpenOffice
    "bs4": "beautifulsoup4",  # Ανάλυση HTML/XML
    "markdown": "Markdown",     # Ανάλυση Markdown
    "lxml": "lxml",           # Ανάλυση XML/HTML
    "cairosvg": "cairosvg",     # Μετατροπή SVG
    "psd_tools": "psd-tools",   # Ανάγνωση PSD
    "striprtf": "striprtf",     # Ανάγνωση RTF
    "EbookLib": "EbookLib",     # Ανάγνωση EPUB
    "pptx": "python-pptx",    # Ανάγνωση PowerPoint
    "rarfile": "rarfile",       # Εξαγωγή RAR
    "py7zr": "py7zr"          # Εξαγωγή 7-Zip
}

# (40) Φάκελοι για δημιουργία
FOLDER_NAMES = [
    # Εικόνες (10)
    "JPG", "PNG", "WEBP", "BMP", "TIFF", "GIF", "ICO", "TGA", "SVG", "PSD",
    # Έγγραφα (12)
    "PDF", "TXT", "DOCX", "ODT", "HTML", "MD", "CSV", "RTF", "EPUB", "JSON", "XML", "PPTX",
    # Αρχεία (5)
    "ZIP", "TAR", "RAR", "7Z", "GZ",
    # Ήχος (7)
    "MP3", "WAV", "OGG", "FLAC", "M4A", "AAC", "WMA",
    # Βίντεο (6)
    "MP4", "MKV", "AVI", "MOV", "WMV", "FLV"
]

# Λίστες βοηθητικών για λογική
IMAGE_FOLDERS = ["JPG", "PNG", "WEBP", "BMP", "TIFF", "GIF", "ICO", "TGA", "SVG", "PSD"]
IMAGE_EXTS = ['.jpg', '.jpeg', '.png', '.webp', '.bmp', '.tiff', '.gif', '.ico', '.tga']
VECTOR_IMAGE_EXTS = ['.svg']
LAYERED_IMAGE_EXTS = ['.psd']
AV_EXTS = ['.mp3', '.wav', '.ogg', '.flac', '.m4a', '.aac', '.wma', '.mp4', '.mkv', '.avi', '.mov', '.wmv', '.flv']
ARCHIVE_EXTS = ['.zip', '.tar', '.gz', '.bz2', '.rar', '.7z']
TEXT_DOC_EXTS = ['.txt', '.docx', '.odt', '.html', '.md', '.csv', '.rtf', '.epub', '.json', '.xml', '.pptx', '.svg']
DATA_EXTS = ['.csv', '.json', '.xml']

# Διαδρομές
STORAGE_PATH = "/storage/emulated/0"
DOWNLOAD_PATH = os.path.join(STORAGE_PATH, "Download")
# --- ΤΡΟΠΟΠΟΙΗΜΕΝΟ: Όλοι οι φάκελοι είναι τώρα μέσα στο "File Converter" ---
BASE_CONVERTER_PATH = os.path.join(DOWNLOAD_PATH, "File Converter")

# Παγκόσμιες σημαίες για εξωτερικά εκτελέσιμα
HAS_FFMPEG = False
HAS_UNRAR = False
HAS_CAIRO = False

# --- 2. ΣΥΝΑΡΤΗΣΕΙΣ ΠΡΟΕΤΟΙΜΑΣΙΑΣ (Τυπική Εκτύπωση) ---

def clear_screen_standard():
    os.system('clear')

def print_header(title):
    """Εκτύπωση μορφοποιημένης κεφαλίδας"""
    print("\n" + "="*60)
    print(f" {title}")
    print("="*60)

def print_status(message, is_error=False):
    """Εκτύπωση μηνύματος κατάστασης"""
    if is_error:
        print(f"❌ ΣΦΑΛΜΑ: {message}")
    else:
        print(f"📢 {message}")

def check_and_install_dependencies():
    """Ελέγχει και εγκαθιστά τις απαιτούμενες βιβλιοθήκες Python."""
    print_header("Έλεγχος Απαιτούμενων Βιβλιοθηκών Python (14)")
    all_installed = True
    for module_name, package_name in REQUIRED_MODULES.items():
        spec = importlib.util.find_spec(module_name)
        if spec is None:
            all_installed = False
            print(f"Εγκατάσταση '{package_name}'...")
            try:
                with open(os.devnull, 'w') as devnull:
                    with redirect_stdout(devnull), redirect_stderr(devnull):
                        subprocess.run([sys.executable, "-m", "pip", "install", package_name], check=True)
                print(f"✅ Εγκαταστάθηκε επιτυχώς το '{package_name}'.")
            except Exception:
                print(f"❌ ΣΦΑΛΜΑ: Αποτυχία εγκατάστασης του '{package_name}'.")
                print(f"Παρακαλώ εγκαταστήστε το χειροκίνητα: pip install {package_name}")
                sys.exit(1)
        else:
            # print(f"Η βιβλιοθήκη '{package_name}' είναι ήδη εγκατεστημένη.")
            pass
    
    if all_installed:
        print("✅ Όλες οι βιβλιοθήκες Python είναι παρόντες.\n")
    else:
        print("✅ Όλες οι απαιτούμενες βιβλιοθήκες είναι τώρα εγκατεστημένες.\n")
    time.sleep(0.5)

def check_external_bins():
    """Ελέγχει για 'ffmpeg', 'unrar', και 'cairo'."""
    global HAS_FFMPEG, HAS_UNRAR, HAS_CAIRO
    print_header("Έλεγχος Εξωτερικών Εκτελέσιμων")
    
    # Έλεγχος ffmpeg
    try:
        with open(os.devnull, 'w') as devnull:
            subprocess.run(["ffmpeg", "-version"], check=True, stdout=devnull, stderr=devnull)
        print("✅ Βρέθηκε 'ffmpeg'. Οι μετατροπές ήχου/βίντεο είναι ΕΝΕΡΓΟΠΟΙΗΜΕΝΕΣ.")
        HAS_FFMPEG = True
    except (subprocess.CalledProcessError, FileNotFoundError):
        print("⚠️  ΠΡΟΕΙΔΟΠΟΙΗΣΗ: Δεν βρέθηκε 'ffmpeg'. Μετατροπές A/V ΑΠΕΝΕΡΓΟΠΟΙΗΜΕΝΕΣ.")
        print("  Για ενεργοποίηση, εκτελέστε: pkg install ffmpeg\n")
        HAS_FFMPEG = False

    # Έλεγχος unrar
    try:
        with open(os.devnull, 'w') as devnull:
            subprocess.run(["unrar"], check=True, stdout=devnull, stderr=devnull)
        print("✅ Βρέθηκε 'unrar'. Η εξαγωγή RAR είναι ΕΝΕΡΓΟΠΟΙΗΜΕΝΗ.")
        HAS_UNRAR = True
    except (subprocess.CalledProcessError, FileNotFoundError):
        print("⚠️  ΠΡΟΕΙΔΟΠΟΙΗΣΗ: Δεν βρέθηκε 'unrar'. Εξαγωγή RAR ΑΠΕΝΕΡΓΟΠΟΙΗΜΕΝΗ.")
        print("  Για ενεργοποίηση, εκτελέστε: pkg install unrar\n")
        HAS_UNRAR = False
        
    # Έλεγχος cairo (για SVG)
    if importlib.util.find_spec("cairosvg") is not None:
        print("✅ Βρέθηκε 'cairosvg'. Οι μετατροπές SVG είναι ΕΝΕΡΓΟΠΟΙΗΜΕΝΕΣ.")
        HAS_CAIRO = True
    else:
        print("⚠️  ΠΡΟΕΙΔΟΠΟΙΗΣΗ: Δεν βρέθηκε βιβλιοθήκη 'cairosvg'. Μετατροπή SVG ΑΠΕΝΕΡΓΟΠΟΙΗΜΕΝΗ.")
        print("  Το script προσπάθησε να την εγκαταστήσει, αλλά μπορεί να απέτυχε.")
        print("  Μπορεί επίσης να χρειαστείτε: pkg install libcairo libgirepository\n")
        HAS_CAIRO = False
        
    print("")
    time.sleep(0.5)


def check_storage_access():
    print_header("Έλεγχος Πρόσβασης Αποθήκευσης")
    if not os.path.exists(DOWNLOAD_PATH):
        print(f"❌ ΣΦΑΛΜΑ: Δεν είναι δυνατή η πρόσβαση στην εσωτερική αποθήκευση στο '{DOWNLOAD_PATH}'.")
        print("Παρακαλώ εκτελέστε 'termux-setup-storage' στο τερματικό του Termux,")
        print("χορήγηση άδειας και μετά εκτελέστε ξανά αυτό το script.")
        sys.exit(1)
    print("✅ Πρόσβαση αποθήκευσης επιβεβαιώθηκε.\n")
    time.sleep(0.5)

def setup_folders():
    # --- ΤΡΟΠΟΠΟΙΗΜΕΝΟ: Δημιουργεί πρώτα τον κύριο φάκελο "File Converter" ---
    print_header("Δημιουργία Φακέλων Οργανωτή")
    print(f"Τοποθεσία: {BASE_CONVERTER_PATH}")
    try:
        # 1. Δημιουργία του κύριου γονικού φακέλου
        os.makedirs(BASE_CONVERTER_PATH, exist_ok=True)
        # 2. Δημιουργία όλων των 40 υπο-φακέλων μέσα σε αυτόν
        for folder in FOLDER_NAMES:
            os.makedirs(os.path.join(BASE_CONVERTER_PATH, folder), exist_ok=True)
        print(f"✅ Δημιουργήθηκαν/επαληθεύτηκαν επιτυχώς {len(FOLDER_NAMES)} υπο-φάκελοι.\n")
    except Exception as e:
        print(f"❌ ΣΦΑΛΜΑ: Δεν ήταν δυνατή η δημιουργία φακέλων: {e}")
        sys.exit(1)
    time.sleep(0.5)

# --- 3. ΕΙΣΑΓΩΓΕΣ (Μετά την Εγκατάσταση) ---
try:
    from PIL import Image
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from docx import Document
    from odf.opendocument import load as odf_load
    from odf.text import P as odf_P
    from bs4 import BeautifulSoup
    import markdown
    import lxml
    import cairosvg
    from psd_tools import PSDImage
    from striprtf.striprtf import rtf_to_text
    from ebooklib import epub, ITEM_DOCUMENT
    import pptx
    import rarfile
    import py7zr
except ImportError as e:
    print(f"❌ ΚΡΙΣΙΜΟ ΣΦΑΛΜΑ: Αποτυχία εισαγωγής βιβλιοθήκης: {e}")
    print("Παρακαλώ βεβαιωθείτε ότι όλες οι εξαρτήσεις είναι εγκατεστημένες (δείτε τα αρχεία καταγραφής εκκίνησης).")
    sys.exit(1)

# --- 4. ΚΥΡΙΑ ΛΟΓΙΚΗ ΜΕΤΑΤΡΟΠΗΣ ---

def get_text_from_file(input_path, in_ext):
    """Βοηθητική για εξαγωγή απλού κειμένου από διάφορους τύπους εγγράφων."""
    text_lines = []
    try:
        if in_ext == '.txt':
            with open(input_path, 'r', encoding='utf-8') as f:
                text_lines = f.readlines()
        elif in_ext == '.docx':
            doc = Document(input_path)
            text_lines = [para.text + '\n' for para in doc.paragraphs]
        elif in_ext == '.odt':
            doc = odf_load(input_path)
            for para in doc.getElementsByType(odf_P):
                text_lines.append(str(para) + '\n')
        elif in_ext in ['.html', '.xml', '.svg']:
            with open(input_path, 'r', encoding='utf-8') as f:
                parser = 'lxml' if in_ext != '.html' else 'html.parser'
                soup = BeautifulSoup(f, parser)
                text_lines = [line + '\n' for line in soup.get_text().splitlines()]
        elif in_ext == '.md':
            with open(input_path, 'r', encoding='utf-8') as f:
                html = markdown.markdown(f.read())
                soup = BeautifulSoup(html, 'html.parser')
                text_lines = [line + '\n' for line in soup.get_text().splitlines()]
        elif in_ext == '.csv':
            with open(input_path, 'r', encoding='utf-8', newline='') as f:
                reader = csv.reader(f)
                text_lines = [','.join(row) + '\n' for row in reader]
        elif in_ext == '.json':
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                text_lines = [json.dumps(data, indent=2)]
        elif in_ext == '.rtf':
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            text_lines = [rtf_to_text(content)]
        elif in_ext == '.epub':
            book = epub.read_epub(input_path)
            for item in book.get_items():
                if item.get_type() == ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text_lines.append(soup.get_text() + '\n\n') # Προσθήκη χώρου μεταξύ κεφαλαίων
        elif in_ext == '.pptx':
            prs = pptx.Presentation(input_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_lines.append(shape.text + '\n')
    except Exception as e:
        raise Exception(f"Αποτυχία εξαγωγής κειμένου: {e}")
    return text_lines

def write_text_to_pdf(text_lines, output_path):
    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    margin_x, margin_y = 0.75 * inch, 1 * inch
    text_object = c.beginText(margin_x, height - margin_y)
    text_object.setFont("Helvetica", 10)
    line_height, y = 12, height - margin_y
    for line in text_lines:
        for sub_line in line.split('\n'): # Χειρισμός συμβολοσειρών πολλαπλών γραμμών
            if y < margin_y:
                c.drawText(text_object)
                c.showPage()
                text_object = c.beginText(margin_x, height - margin_y)
                text_object.setFont("Helvetica", 10)
                y = height - margin_y
            text_object.textLine(sub_line.strip('\r'))
            y -= line_height
    c.drawText(text_object)
    c.save()

def handle_image_conversion(in_path, out_path):
    # (Χειριστής Pillow, αμετάβλητος)
    with Image.open(in_path) as img:
        if out_path.lower().endswith(('.jpg', '.jpeg')):
            if img.mode == 'RGBA':
                img = img.convert('RGB')
        img.save(out_path)

def handle_svg_conversion(in_path, out_path):
    """Μετατρέπει SVG σε PNG ή PDF."""
    if not HAS_CAIRO:
        raise Exception("Βιβλιοθήκες Cairo/SVG δεν είναι εγκατεστημένες.")
    out_ext = os.path.splitext(out_path)[1].lower()
    if out_ext == '.png':
        cairosvg.svg2png(url=in_path, write_to=out_path)
    elif out_ext == '.pdf':
        cairosvg.svg2pdf(url=in_path, write_to=out_path)
    else:
        raise Exception(f"Μετατροπή SVG σε {out_ext} δεν υποστηρίζεται.")

def handle_psd_conversion(in_path, out_path):
    """Μετατρέπει σύνθεση PSD σε επίπεδη εικόνα."""
    psd = PSDImage.open(in_path)
    composite_image = psd.composite()
    composite_image.save(out_path)

def handle_av_conversion(in_path, out_path):
    # (Αμετάβλητος)
    if not HAS_FFMPEG:
        raise Exception("Δεν βρέθηκε 'ffmpeg'. Μετατροπή A/V είναι απενεργοποιημένη.")
    command = ['ffmpeg', '-i', in_path, '-y', out_path]
    print_header("Εκτέλεση ffmpeg")
    print(f"Εντολή: {' '.join(command)}")
    print("Αυτό μπορεί να πάρει λίγο χρόνο...")
    try:
        with open(os.devnull, 'w') as devnull:
            subprocess.run(command, check=True, stdout=devnull, stderr=subprocess.STDOUT)
        print("✅ Μετατροπή ffmpeg ολοκληρώθηκε επιτυχώς.")
    except Exception as e:
        print(f"❌ ΣΦΑΛΜΑ ffmpeg: {e}")
        raise Exception(f"Αποτυχία μετατροπής ffmpeg. {e}")

def handle_extraction(in_path, out_folder_path, in_ext):
    """Εξάγει διάφορους τύπους αρχείων."""
    base_name = os.path.splitext(os.path.basename(in_path))[0]
    extract_path = os.path.join(out_folder_path, base_name)
    os.makedirs(extract_path, exist_ok=True)
    
    if in_ext == '.zip':
        with zipfile.ZipFile(in_path, 'r') as zf:
            zf.extractall(extract_path)
    elif in_ext in ['.tar', '.gz', '.bz2']:
        if in_ext == '.gz' and not in_path.endswith('.tar.gz'): # Μονού αρχείου gzip
             out_filename = os.path.splitext(os.path.basename(in_path))[0]
             out_path = os.path.join(out_folder_path, out_filename) # Εξαγωγή σε φάκελο, όχι υποφάκελο
             with gzip.open(in_path, 'rb') as f_in:
                 with open(out_path, 'wb') as f_out:
                     shutil.copyfileobj(f_in, f_out)
             return f"Αποσυμπιέστηκε σε: {out_path}" # Διαφορετικό μήνυμα
        else: # .tar, .tar.gz, .tar.bz2
            with tarfile.open(in_path, 'r:*') as tf:
                tf.extractall(extract_path)
    elif in_ext == '.rar':
        if not HAS_UNRAR:
            raise Exception("Δεν βρέθηκε εκτελέσιμο 'unrar'.")
        with rarfile.RarFile(in_path) as rf:
            rf.extractall(extract_path)
    elif in_ext == '.7z':
        with py7zr.SevenZipFile(in_path, 'r') as zf:
            zf.extractall(extract_path)
            
    return f"Εξήχθη σε: {extract_path}" # Προεπιλεγμένο μήνυμα επιτυχίας

def handle_data_conversion(in_path, out_path, in_ext, out_ext):
    """Χειρίζεται μετατροπές CSV <-> JSON."""
    if in_ext == '.csv' and out_ext == '.json':
        data = []
        with open(in_path, 'r', encoding='utf-8', newline='') as f:
            reader = csv.DictReader(f)
            for row in reader:
                data.append(row)
        with open(out_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=4)
    elif in_ext == '.json' and out_ext == '.csv':
        with open(in_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        if not isinstance(data, list) or not data:
            raise Exception("Το JSON πρέπει να είναι μια μη κενή λίστα αντικειμένων.")
        if not all(isinstance(x, dict) for x in data):
            raise Exception("Το JSON πρέπει να είναι μια λίστα αντικειμένων (λεξικά).")
            
        headers = data[0].keys()
        with open(out_path, 'w', encoding='utf-8', newline='') as f:
            writer = csv.DictWriter(f, fieldnames=headers)
            writer.writeheader()
            writer.writerows(data)
    else:
        raise Exception(f"Μετατροπή δεδομένων {in_ext} σε {out_ext} δεν υποστηρίζεται.")

def handle_md_to_html(in_path, out_path):
    # (Αμετάβλητος)
    with open(in_path, 'r', encoding='utf-8') as f:
        html = markdown.markdown(f.read())
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(html)

def handle_single_image_to_pdf(in_path, out_path):
    # (Αμετάβλητος)
    try:
        with Image.open(in_path) as img:
            img_rgb = img.convert('RGB')
            img_rgb.save(out_path, "PDF", resolution=100.0)
    except Exception as e:
        raise Exception(f"Σφάλμα Pillow (Εικόνα->PDF): {e}")

def handle_multi_image_to_pdf(image_paths, out_path):
    # (Αμετάβλητος)
    try:
        images_rgb = []
        for path in image_paths:
            img = Image.open(path)
            images_rgb.append(img.convert('RGB'))
        if not images_rgb:
            raise Exception("Δεν βρέθηκαν εικόνες για μετατροπή.")
        images_rgb[0].save(
            out_path, "PDF", resolution=100.0,
            save_all=True, append_images=images_rgb[1:]
        )
    except Exception as e:
        raise Exception(f"Σφάλμα Pillow (Πολλαπλών Εικόνων->PDF): {e}")

# --- 5. ΚΥΡΙΟΣ ΔΙΑΚΟΜΙΣΤΗΣ ΜΕΤΑΤΡΟΠΩΝ ---

def convert_file(in_path, out_folder_name):
    """
    Κύρια συνάρτηση δρομολόγησης για την αποστολή εργασιών μετατροπής.
    Επιστρέφει (επιτυχία_boolean, μήνυμα_συμβολοσειράς)
    """
    in_ext = os.path.splitext(in_path)[1].lower()
    out_ext = f".{out_folder_name.lower()}"
    
    base_name = os.path.splitext(os.path.basename(in_path))[0]
    out_folder_path = os.path.join(BASE_CONVERTER_PATH, out_folder_name)
    out_path = os.path.join(out_folder_path, f"{base_name}{out_ext}")

    try:
        # --- Διαδρομή 1: Εξαγωγή ---
        if in_ext in ARCHIVE_EXTS:
            # Σημείωση: Τα αρχεία GZ θα αποσυμπιεστούν *μέσα* στο φάκελο με όνομα GZ
            # Όλα τα άλλα (ZIP, TAR, RAR, 7Z) εξάγονται σε *υπο-φάκελο*
            out_folder = out_folder_path if in_ext == '.gz' else os.path.join(BASE_CONVERTER_PATH, out_folder_name)
            message = handle_extraction(in_path, out_folder, in_ext)
            return (True, message)

        # --- Διαδρομή 2: Μετατροπή SVG (σε PNG, PDF) ---
        if in_ext in VECTOR_IMAGE_EXTS and out_ext in ['.png', '.pdf']:
            handle_svg_conversion(in_path, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")

        # --- Διαδρομή 3: Μετατροπή PSD (σε επίπεδη εικόνα) ---
        if in_ext in LAYERED_IMAGE_EXTS and out_ext in IMAGE_EXTS:
            handle_psd_conversion(in_path, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")

        # --- Διαδρομή 4: Εικόνα-σε-Εικόνα (Pillow) ---
        if in_ext in IMAGE_EXTS and out_ext in IMAGE_EXTS:
            handle_image_conversion(in_path, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")
            
        # --- Διαδρομή 5: Μονής Εικόνας-σε-PDF ---
        if in_ext in IMAGE_EXTS and out_ext == '.pdf':
            handle_single_image_to_pdf(in_path, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")

        # --- Διαδρομή 6: A/V-σε-A/V (ffmpeg) ---
        if in_ext in AV_EXTS and out_ext in AV_EXTS:
            handle_av_conversion(in_path, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")
            
        # --- Διαδρομή 7: Μετατροπή Δεδομένων (CSV <-> JSON) ---
        if in_ext in ['.csv', '.json'] and out_ext in ['.csv', '.json']:
            handle_data_conversion(in_path, out_path, in_ext, out_ext)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")

        # --- Διαδρομή 8: MD-σε-HTML ---
        if in_ext == '.md' and out_ext == '.html':
            handle_md_to_html(in_path, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")

        # --- Διαδρομή 9: Οτιδήποτε-σε-TXT ---
        if out_ext == '.txt' and in_ext in TEXT_DOC_EXTS:
            text_lines = get_text_from_file(in_path, in_ext)
            with open(out_path, 'w', encoding='utf-8') as f:
                f.writelines(text_lines)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")
            
        # --- Διαδρομή 10: Οτιδήποτε-σε-PDF ---
        if out_ext == '.pdf' and in_ext in TEXT_DOC_EXTS:
            text_lines = get_text_from_file(in_path, in_ext)
            write_text_to_pdf(text_lines, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")

        # --- Δεν Βρέθηκε Διαδρομή ---
        return (False, f"Μη υποστηριζόμενη μετατροπή: {in_ext} σε {out_ext}")

    except Exception as e:
        return (False, f"ΣΦΑΛΜΑ: {str(e)}")

# --- 6. ΣΥΣΤΗΜΑ ΜΕΝΟΥ ΜΕ ΑΡΙΘΜΗΣΗ ---

def run_menu(title, options, sub_title=""):
    """Εμφάνιση αριθμημένου μενού και λήψη επιλογής χρήστη"""
    while True:
        clear_screen_standard()
        print_header(title)
        if sub_title:
            print(f"\n{sub_title}\n")
        
        for i, option in enumerate(options, 1):
            print(f"{i:2d}. {option}")
        
        print(f"\n 0. Πίσω")
        
        try:
            choice = int(input("\nΕισάγετε την επιλογή σας (αριθμός): "))
            if choice == 0:
                return None
            if 1 <= choice <= len(options):
                return options[choice - 1]
            else:
                print_status("Μη έγκυρη επιλογή. Παρακαλώ δοκιμάστε ξανά.", is_error=True)
                input("Πατήστε Enter για συνέχεια...")
        except ValueError:
            print_status("Παρακαλώ εισάγετε έγκυρο αριθμό.", is_error=True)
            input("Πατήστε Enter για συνέχεια...")

def run_file_selector(folder_path, title, input_folder_name):
    """Εμφάνιση αρχείων σε φάκελο για επιλογή"""
    try:
        files = [f for f in os.listdir(folder_path) if os.path.isfile(os.path.join(folder_path, f))]
        files.sort()
    except Exception as e:
        print_status(f"Σφάλμα ανάγνωσης {folder_path}: {e}", is_error=True)
        input("Πατήστε Enter για συνέχεια...")
        return None
    
    if not files:
        print_status(f"Δεν βρέθηκαν αρχεία στον {os.path.basename(folder_path)}", is_error=True)
        input("Πατήστε Enter για συνέχεια...")
        return None
        
    options = ["[ .. Επιστροφή .. ]"]
    
    if input_folder_name in IMAGE_FOLDERS:
        options.append(f"[ ** Μετατροπή ΟΛΩΝ {len(files)} Εικόνων στον '{input_folder_name}' σε ένα PDF ** ]")
    
    options.extend(files)
    
    selection = run_menu(title, options, f"Φάκελος: /Download/File Converter/{input_folder_name}")
    if selection == "[ .. Επιστροφή .. ]":
        return None
    return selection

def run_confirmation(prompt):
    """Λήψη επιβεβαίωσης από χρήστη"""
    options = ["Ναι", "Όχι"]
    selection = run_menu("Επιβεβαίωση", options, prompt)
    return selection

def run_help():
    """Εμφάνιση πληροφοριών βοήθειας"""
    clear_screen_standard()
    print_header("Πώς να Χρησιμοποιήσετε")
    help_text = [
        "Αυτός ο μετατροπέας χρησιμοποιεί μια απλή διαδικασία 3 βημάτων:",
        "",
        "1. ΜΕΤΑΚΙΝΗΣΤΕ ΤΑ ΑΡΧΕΙΑ ΣΑΣ:",
        "   Χρησιμοποιήστε τη Διαχείριση Αρχείων του τηλεφώνου σας. Μεταβείτε στο:",
        f"   /Download/File Converter/",
        "   Μετακινήστε αρχεία στον σωστό φάκελο (π.χ., μετακινήστε",
        "   το 'report.docx' στον φάκελο 'DOCX').",
        "",
        "2. ΕΚΤΕΛΕΣΤΕ ΑΥΤΟΝ ΤΟΝ ΜΕΤΑΤΡΟΠΕΑ:",
        "   Επιλέξτε 'Μετατροπή Αρχείου' από το κύριο μενού.",
        "",
        "3. ΑΚΟΛΟΥΘΗΣΤΕ ΤΑ ΒΗΜΑΤΑ:",
        "   Βήμα 1: Επιλέξτε τον φάκελο ΕΙΣΟΔΟΥ (π.χ., 'DOCX').",
        "   Βήμα 2: Επιλέξτε το αρχείο που θέλετε να μετατρέψετε.",
        "   Βήμα 3: Επιλέξτε τη ΜΟΡΦΗ ΕΞΟΔΟΥ (π.χ., 'PDF').",
        "",
        "** ΕΙΔΙΚΕΣ ΜΕΤΑΤΡΟΠΕΣ **",
        " - Αρχεία (ZIP, RAR, 7Z, TAR): Επιλέξτε 'ZIP' -> 'file.zip' -> 'ZIP'",
        "   Αυτό θα εξάγει το 'file.zip' σε έναν νέο φάκελο: /ZIP/file/",
        " - PDF Πολλαπλών Εικόνων: Επιλέξτε 'JPG' -> '[ ** Μετατροπή ΟΛΩΝ... ** ]'",
        "   Αυτό συνδυάζει όλες τις εικόνες στον 'JPG' σε ένα PDF.",
        " - Δεδομένα: Μπορείτε να μετατρέψετε CSV <-> JSON.",
        " - A/V: Οι μετατροπές ήχου/βίντεο απαιτούν 'ffmpeg' (δείτε την εκκίνηση)."
    ]
    
    for line in help_text:
        print(line)
    
    print("\n" + "="*60)
    input("Πατήστε Enter για επιστροφή...")

def run_text_input(prompt):
    """Λήψη κειμένου από χρήστη"""
    clear_screen_standard()
    print_header("Απαιτείται Εισαγωγή")
    print(f"\n{prompt}")
    print("\nΠληκτρολογήστε το όνομα αρχείου (χωρίς επέκταση). Πατήστε Enter όταν τελειώσετε.")
    print("Πληκτρολογήστε 'q' για ακύρωση.")
    
    text = input("\nΕισάγετε όνομα αρχείου: ").strip()
    if text.lower() == 'q':
        return None
    return text

def run_multi_image_to_pdf_wizard(input_folder_path, input_folder_name):
    """Χειρισμός μετατροπής πολλαπλών εικόνων σε PDF"""
    try:
        image_paths = [
            os.path.join(input_folder_path, f) 
            for f in os.listdir(input_folder_path) 
            if os.path.splitext(f)[1].lower() in IMAGE_EXTS
        ]
        image_paths.sort()
    except Exception as e:
        print_status(f"Σφάλμα ανάγνωσης εικόνων: {e}", is_error=True)
        input("Πατήστε Enter για συνέχεια...")
        return
        
    if not image_paths:
        print_status("Δεν βρέθηκαν εικόνες σε αυτόν τον φάκελο.", is_error=True)
        input("Πατήστε Enter για συνέχεια...")
        return
        
    confirm = run_confirmation(f"Συνδυασμός και των {len(image_paths)} εικόνων στον '{input_folder_name}' σε ένα PDF;")
    if confirm != "Ναι":
        return
        
    default_name = f"{input_folder_name}_Άλμπουμ"
    filename = run_text_input(f"Εισάγετε ένα όνομα για το PDF (προεπιλογή: {default_name})")
    if filename is None:
        return
    if not filename:
        filename = default_name
        
    out_folder_path = os.path.join(BASE_CONVERTER_PATH, "PDF")
    out_path = os.path.join(out_folder_path, f"{filename}.pdf")
    
    print_status("Εργασία... συνδυασμός εικόνων σε PDF...")
    try:
        handle_multi_image_to_pdf(image_paths, out_path)
        print_status(f"Επιτυχία! Αποθηκεύτηκε σε: /PDF/{filename}.pdf")
    except Exception as e:
        print_status(f"ΣΦΑΛΜΑ: {e}", is_error=True)
    input("Πατήστε Enter για συνέχεια...")

# --- 7. ΚΥΡΙΑ ΕΦΑΡΜΟΓΗ ---

def main():
    """Κύριος βρόγχος εφαρμογής"""
    while True:
        main_choice = run_menu("Κύριο Μενού", ["Μετατροπή Αρχείου", "Βοήθεια / Πώς να Χρησιμοποιήσετε", "Έξοδος"])
        if main_choice == "Έξοδος" or main_choice is None:
            break
        if main_choice == "Βοήθεια / Πώς να Χρησιμοποιήσετε":
            run_help()
            continue
            
        # Βήμα 1: Επιλογή φακέλου εισόδου
        input_folder = run_menu("Βήμα 1: Επιλογή Φακέλου ΕΙΣΟΔΟΥ", FOLDER_NAMES)
        if input_folder is None:
            continue
            
        input_folder_path = os.path.join(BASE_CONVERTER_PATH, input_folder)
        
        # Βήμα 2: Επιλογή αρχείου
        input_file = run_file_selector(input_folder_path, f"Βήμα 2: Επιλογή Αρχείου από '{input_folder}'", input_folder)
        if input_file is None:
            continue
        
        if input_file.startswith("[ ** Μετατροπή ΟΛΩΝ"):
            run_multi_image_to_pdf_wizard(input_folder_path, input_folder)
            continue
            
        full_input_path = os.path.join(input_folder_path, input_file)
        
        # --- Ειδική Περίπτωση: Εξαγωγή Αρχείου ---
        in_ext = os.path.splitext(input_file)[1].lower()
        if in_ext in ARCHIVE_EXTS:
            # Για αρχεία, ο "φάκελος εξόδου" είναι απλώς ο φάκελος του ίδιου τύπου
            output_folder = input_folder
            prompt = f"Εξαγωγή '{input_file}' στον '/{output_folder}/';"
            if in_ext not in ['.zip', '.rar', '.7z', '.tar', '.gz', '.bz2']:
                output_folder = input_folder # Ασφάλεια
                
        else:
            # --- Κανονική Διαδρομή Μετατροπής ---
            output_folder = run_menu("Βήμα 3: Επιλογή Μορφής/Φακέλου ΕΞΟΔΟΥ", FOLDER_NAMES)
            if output_folder is None:
                continue
            if output_folder == input_folder:
                print_status("Σφάλμα: Ο φάκελος εισόδου και εξόδου δεν μπορεί να είναι ο ίδιος.", is_error=True)
                input("Πατήστε Enter για συνέχεια...")
                continue
            prompt = f"Μετατροπή '{input_file}' σε μορφή {output_folder};"
             
        confirm = run_confirmation(prompt)
        if confirm != "Ναι":
            continue

        print_status("Εργασία, παρακαλώ περιμένετε...")
        success, message = convert_file(full_input_path, output_folder)
        print_status(message, is_error=not success)
        input("Πατήστε Enter για συνέχεια...")

# --- 8. ΣΗΜΕΙΟ ΕΝΑΡΞΗΣ SCRIPT ---

if __name__ == "__main__":
    try:
        clear_screen_standard()
        print("--- Αρχικοποίηση Termux Converter v3.1 (40 Μορφές) ---")
        check_and_install_dependencies()
        check_external_bins()
        check_storage_access()
        setup_folders()
        
        print("--- Ρύθμιση Ολοκληρώθηκε ---")
        # --- ΤΡΟΠΟΠΟΙΗΜΕΝΟ: Ενημερωμένο μήνυμα τελικής διαδρομής ---
        print(f"Οι φάκελοι οργανωτή είναι έτοιμοι στο: /storage/emulated/0/Download/File Converter/")
        print("\nΕκκίνηση εφαρμογής...")
        time.sleep(1)
        
        main()
        print("Ο Μετατροπέας Αρχείων τερμάτισε επιτυχώς.")

    except KeyboardInterrupt:
        print("\nΈξοδος...")
    except Exception as e:
        print("\nΠροέκυψε ένα κρίσιμο σφάλμα:")
        traceback.print_exc()
    finally:
        os.system('clear')