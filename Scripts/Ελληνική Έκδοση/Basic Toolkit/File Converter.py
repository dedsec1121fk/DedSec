#!/usr/bin/env python

import sys
import os
import subprocess
import importlib.util
import time
import curses
import traceback
import zipfile
import tarfile
import csv
import json
import gzip
import shutil
from contextlib import redirect_stderr, redirect_stdout

# --- 1. ΡΥΘΜΙΣΕΙΣ & ΔΙΑΜΟΡΦΩΣΗ ---

# Απαιτούμενα πακέτα συστήματος για Termux (Δυαδικά + Εργαλεία κατασκευής για βιβλιοθήκες Python)
TERMUX_PACKAGES = [
    "ffmpeg", 
    "unrar", 
    "libcairo", 
    "libgirepository", 
    "libjpeg-turbo", 
    "libpng", 
    "python", 
    "clang",       # Απαιτείται για την κατασκευή ορισμένων βιβλιοθηκών python
    "make",        # Απαιτείται για την κατασκευή ορισμένων βιβλιοθηκών python
    "binutils",    # Απαιτείται για την κατασκευή ορισμένων βιβλιοθηκών python
    "libffi"
]

# (14) Βιβλιοθήκες Python για αυτόματη εγκατάσταση
REQUIRED_MODULES = {
    "Pillow": "Pillow",         # Εικόνες
    "reportlab": "reportlab",   # Δημιουργία PDF
    "docx": "python-docx",    # Έγγραφα Word
    "odf": "odfpy",           # Έγγραφα OpenOffice
    "bs4": "beautifulsoup4",  # Ανάλυση HTML/XML
    "markdown": "Markdown",     # Ανάλυση Markdown
    "lxml": "lxml",           # Ανάλυση XML/HTML
    "cairosvg": "cairosvg",     # Μετατροπή SVG
    "psd_tools": "psd-tools",   # Ανάγνωση PSD
    "striprtf": "striprtf",     # Ανάγνωση RTF
    "EbookLib": "EbookLib",     # Ανάγνωση EPUB
    "pptx": "python-pptx",    # Ανάγνωση PowerPoint
    "rarfile": "rarfile",       # Εξαγωγή RAR
    "py7zr": "py7zr"          # Εξαγωγή 7-Zip
}

# (40) Φάκελοι για δημιουργία
FOLDER_NAMES = [
    "JPG", "PNG", "WEBP", "BMP", "TIFF", "GIF", "ICO", "TGA", "SVG", "PSD",
    "PDF", "TXT", "DOCX", "ODT", "HTML", "MD", "CSV", "RTF", "EPUB", "JSON", "XML", "PPTX",
    "ZIP", "TAR", "RAR", "7Z", "GZ",
    "MP3", "WAV", "OGG", "FLAC", "M4A", "AAC", "WMA",
    "MP4", "MKV", "AVI", "MOV", "WMV", "FLV"
]

IMAGE_FOLDERS = ["JPG", "PNG", "WEBP", "BMP", "TIFF", "GIF", "ICO", "TGA", "SVG", "PSD"]
IMAGE_EXTS = ['.jpg', '.jpeg', '.png', '.webp', '.bmp', '.tiff', '.gif', '.ico', '.tga']
VECTOR_IMAGE_EXTS = ['.svg']
LAYERED_IMAGE_EXTS = ['.psd']
AV_EXTS = ['.mp3', '.wav', '.ogg', '.flac', '.m4a', '.aac', '.wma', '.mp4', '.mkv', '.avi', '.mov', '.wmv', '.flv']
ARCHIVE_EXTS = ['.zip', '.tar', '.gz', '.bz2', '.rar', '.7z']
TEXT_DOC_EXTS = ['.txt', '.docx', '.odt', '.html', '.md', '.csv', '.rtf', '.epub', '.json', '.xml', '.pptx', '.svg']

STORAGE_PATH = "/storage/emulated/0"
DOWNLOAD_PATH = os.path.join(STORAGE_PATH, "Download")
BASE_CONVERTER_PATH = os.path.join(DOWNLOAD_PATH, "File Converter")

HAS_FFMPEG = False
HAS_UNRAR = False
HAS_CAIRO = False

# --- 2. ΣΥΝΑΡΤΗΣΕΙΣ ΑΥΤΟΜΑΤΗΣ ΕΓΚΑΤΑΣΤΑΣΗΣ ---

def clear_screen_standard():
    os.system('clear')

def install_termux_system_deps():
    """Εγκαθιστά αυτόματα τα απαιτούμενα πακέτα συστήματος μέσω pkg."""
    print("--- 1/4 Έλεγχος Εξαρτήσεων Συστήματος (Termux) ---")
    
    # Έλεγχος αν τρέχουμε πιθανώς σε Termux
    if not os.path.exists("/data/data/com.termux/files/usr/bin/pkg"):
        print("Δεν τρέχει σε τυπικό περιβάλλον Termux. Παράκαμψη εγκατάστασης πακέτων συστήματος.")
        return

    print("Ενημέρωση λιστών πακέτων και εγκατάσταση δυαδικών...")
    print("Αυτό μπορεί να πάρει λίγο χρόνο. Παρακαλώ περιμένετε...")
    
    try:
        # Ενημέρωση αποθετηρίων
        subprocess.run(["pkg", "update", "-y"], check=False)
        
        # Εγκατάσταση πακέτων
        cmd = ["pkg", "install", "-y"] + TERMUX_PACKAGES
        subprocess.run(cmd, check=True)
        print("Τα δυαδικά συστήματος εγκαταστάθηκαν με επιτυχία.\n")
    except Exception as e:
        print(f"ΠΡΟΕΙΔΟΠΟΙΗΣΗ: Η εγκατάσταση συστήματος απέτυχε: {e}")
        print("Προσπάθεια συνέχισης, αλλά ορισμένες λειτουργίες (FFmpeg, Cairo) μπορεί να αποτύχουν.\n")
    time.sleep(1)

def check_and_install_python_deps():
    """Ελέγχει και εγκαθιστά τα απαιτούμενα modules Python."""
    print("--- 2/4 Έλεγχος Βιβλιοθηκών Python ---")
    
    # Αναβάθμιση pip πρώτα για να διασφαλιστεί ότι μπορούμε να εγκαταστήσουμε σωστά τα wheels
    try:
        subprocess.run([sys.executable, "-m", "pip", "install", "--upgrade", "pip"], 
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except:
        pass

    all_installed = True
    for module_name, package_name in REQUIRED_MODULES.items():
        spec = importlib.util.find_spec(module_name)
        if spec is None:
            all_installed = False
            print(f"Εγκατάσταση '{package_name}'... (Η κατασκευή wheels μπορεί να πάρει χρόνο)")
            try:
                # Χρησιμοποιούμε απευθείας subprocess για να δούμε την έξοδο αν αποτύχει,
                # αλλά την κρύβουμε αν επιτύχει για να κρατήσουμε το UI καθαρό.
                subprocess.run([sys.executable, "-m", "pip", "install", package_name], check=True)
                print(f"✔ Εγκαταστάθηκε το '{package_name}'.")
            except Exception:
                print(f"✖ ΣΦΑΛΜΑ: Αποτυχία εγκατάστασης του '{package_name}'.")
                print(f"  Δοκιμάστε χειροκίνητα: pip install {package_name}")
                # Δεν τερματίζουμε εδώ, προσπαθούμε να συνεχίσουμε
        else:
            pass
    
    if all_installed:
        print("Όλες οι βιβλιοθήκες Python είναι παρόντες.\n")
    else:
        print("Ο έλεγχος βιβλιοθηκών ολοκληρώθηκε.\n")
    time.sleep(0.5)

def check_external_bins_status():
    """Ελέγχει ποια εξωτερικά δυαδικά είναι πραγματικά διαθέσιμα μετά την εγκατάσταση."""
    global HAS_FFMPEG, HAS_UNRAR, HAS_CAIRO
    print("--- 3/4 Επαλήθευση Ολοκληρωμάτων ---")
    
    # Έλεγχος ffmpeg
    if shutil.which("ffmpeg"):
        print("✔ Βρέθηκε το 'ffmpeg'. Μετατροπές ήχου/βίντεο ΕΝΕΡΓΟΠΟΙΗΜΕΝΕΣ.")
        HAS_FFMPEG = True
    else:
        print("✖ Δεν βρέθηκε το 'ffmpeg'. Μετατροπές ήχου/βίντεο ΑΠΕΝΕΡΓΟΠΟΙΗΜΕΝΕΣ.")
        HAS_FFMPEG = False

    # Έλεγχος unrar
    if shutil.which("unrar"):
        print("✔ Βρέθηκε το 'unrar'. Εξαγωγή RAR ΕΝΕΡΓΟΠΟΙΗΜΕΝΗ.")
        HAS_UNRAR = True
    else:
        print("✖ Δεν βρέθηκε το 'unrar'. Εξαγωγή RAR ΑΠΕΝΕΡΓΟΠΟΙΗΜΕΝΗ.")
        HAS_UNRAR = False
        
    # Έλεγχος cairo
    if importlib.util.find_spec("cairosvg") is not None:
        # Βασικός έλεγχος αν η βιβλιοθήκη φορτώνεται χωρίς σφάλματα DLL
        try:
            import cairosvg
            HAS_CAIRO = True
            print("✔ Φορτώθηκε το 'cairosvg'. Μετατροπές SVG ΕΝΕΡΓΟΠΟΙΗΜΕΝΕΣ.")
        except OSError:
            print("✖ Το 'cairosvg' είναι εγκατεστημένο αλλά λείπουν βιβλιοθήκες συστήματος. SVG ΑΠΕΝΕΡΓΟΠΟΙΗΜΕΝΟ.")
            HAS_CAIRO = False
    else:
        print("✖ Δεν βρέθηκε το 'cairosvg'. SVG ΑΠΕΝΕΡΓΟΠΟΙΗΜΕΝΟ.")
        HAS_CAIRO = False
        
    print("")
    time.sleep(0.5)

def ensure_storage_access():
    print("--- 4/4 Έλεγχος Πρόσβασης Αποθήκευσης ---")
    if not os.path.exists(DOWNLOAD_PATH):
        print(f"Δεν επιτρέπεται η πρόσβαση στο '{DOWNLOAD_PATH}' ή λείπει.")
        print("Προσπάθεια αίτησης δικαιωμάτων αποθήκευσης...")
        try:
            subprocess.run(["termux-setup-storage"], check=True)
            print("Ζητήθηκε άδεια. Παρακαλώ επιτρέψτε την στο popup.")
            print("Αναμονή 5 δευτερολέπτων για διάδοση της άδειας...")
            time.sleep(5)
        except FileNotFoundError:
            print("Δεν ήταν δυνατή η εκτέλεση του 'termux-setup-storage'.")
        
        if not os.path.exists(DOWNLOAD_PATH):
             print(f"ΣΦΑΛΜΑ: Ακόμη δεν είναι δυνατή η πρόσβαση στο '{DOWNLOAD_PATH}'.")
             print("Παρακαλώ επανεκκινήστε το Termux και ξανατρέξτε το script.")
             sys.exit(1)
             
    print("Επιβεβαιώθηκε η πρόσβαση αποθήκευσης.\n")
    time.sleep(0.5)

def setup_folders():
    try:
        os.makedirs(BASE_CONVERTER_PATH, exist_ok=True)
        for folder in FOLDER_NAMES:
            os.makedirs(os.path.join(BASE_CONVERTER_PATH, folder), exist_ok=True)
    except Exception as e:
        print(f"ΣΦΑΛΜΑ: Δεν ήταν δυνατή η δημιουργία φακέλων: {e}")
        sys.exit(1)

# --- 3. ΕΙΣΑΓΩΓΕΣ (Μετά την εγκατάσταση) ---
# Αυτές εισάγονται μέσα σε συναρτήσεις ή τυλίγονται για να αποφευχθεί crash πριν ολοκληρωθεί η εγκατάσταση

# --- 4. ΒΑΣΙΚΗ ΛΟΓΙΚΗ ΜΕΤΑΤΡΟΠΗΣ ---

def get_text_from_file(input_path, in_ext):
    text_lines = []
    try:
        if in_ext == '.txt':
            with open(input_path, 'r', encoding='utf-8') as f:
                text_lines = f.readlines()
        elif in_ext == '.docx':
            from docx import Document
            doc = Document(input_path)
            text_lines = [para.text + '\n' for para in doc.paragraphs]
        elif in_ext == '.odt':
            from odf.opendocument import load as odf_load
            from odf.text import P as odf_P
            doc = odf_load(input_path)
            for para in doc.getElementsByType(odf_P):
                text_lines.append(str(para) + '\n')
        elif in_ext in ['.html', '.xml', '.svg']:
            from bs4 import BeautifulSoup
            with open(input_path, 'r', encoding='utf-8') as f:
                parser = 'lxml' if in_ext != '.html' else 'html.parser'
                soup = BeautifulSoup(f, parser)
                text_lines = [line + '\n' for line in soup.get_text().splitlines()]
        elif in_ext == '.md':
            import markdown
            from bs4 import BeautifulSoup
            with open(input_path, 'r', encoding='utf-8') as f:
                html = markdown.markdown(f.read())
                soup = BeautifulSoup(html, 'html.parser')
                text_lines = [line + '\n' for line in soup.get_text().splitlines()]
        elif in_ext == '.csv':
            with open(input_path, 'r', encoding='utf-8', newline='') as f:
                reader = csv.reader(f)
                text_lines = [','.join(row) + '\n' for row in reader]
        elif in_ext == '.json':
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                text_lines = [json.dumps(data, indent=2)]
        elif in_ext == '.rtf':
            from striprtf.striprtf import rtf_to_text
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            text_lines = [rtf_to_text(content)]
        elif in_ext == '.epub':
            from ebooklib import epub, ITEM_DOCUMENT
            from bs4 import BeautifulSoup
            book = epub.read_epub(input_path)
            for item in book.get_items():
                if item.get_type() == ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text_lines.append(soup.get_text() + '\n\n')
        elif in_ext == '.pptx':
            import pptx
            prs = pptx.Presentation(input_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_lines.append(shape.text + '\n')
    except Exception as e:
        raise Exception(f"Αποτυχία εξαγωγής κειμένου: {e}")
    return text_lines

def write_text_to_pdf(text_lines, output_path):
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    
    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    margin_x, margin_y = 0.75 * inch, 1 * inch
    text_object = c.beginText(margin_x, height - margin_y)
    text_object.setFont("Helvetica", 10)
    line_height, y = 12, height - margin_y
    for line in text_lines:
        for sub_line in line.split('\n'):
            if y < margin_y:
                c.drawText(text_object)
                c.showPage()
                text_object = c.beginText(margin_x, height - margin_y)
                text_object.setFont("Helvetica", 10)
                y = height - margin_y
            text_object.textLine(sub_line.strip('\r'))
            y -= line_height
    c.drawText(text_object)
    c.save()

def handle_image_conversion(stdscr, in_path, out_path):
    from PIL import Image
    with Image.open(in_path) as img:
        if out_path.lower().endswith(('.jpg', '.jpeg')):
            if img.mode == 'RGBA':
                img = img.convert('RGB')
        img.save(out_path)

def handle_svg_conversion(stdscr, in_path, out_path):
    if not HAS_CAIRO:
        raise Exception("Οι βιβλιοθήκες Cairo/SVG δεν είναι εγκατεστημένες/φορτωμένες.")
    import cairosvg
    out_ext = os.path.splitext(out_path)[1].lower()
    if out_ext == '.png':
        cairosvg.svg2png(url=in_path, write_to=out_path)
    elif out_ext == '.pdf':
        cairosvg.svg2pdf(url=in_path, write_to=out_path)
    else:
        raise Exception(f"Μετατροπή SVG σε {out_ext} δεν υποστηρίζεται.")

def handle_psd_conversion(stdscr, in_path, out_path):
    from psd_tools import PSDImage
    psd = PSDImage.open(in_path)
    composite_image = psd.composite()
    composite_image.save(out_path)

def handle_av_conversion(stdscr, in_path, out_path):
    if not HAS_FFMPEG:
        raise Exception("Δεν βρέθηκε το 'ffmpeg'. Η μετατροπή ήχου/βίντεο είναι απενεργοποιημένη.")
    command = ['ffmpeg', '-i', in_path, '-y', out_path]
    curses.endwin()
    print("--- Εκτέλεση ffmpeg ---")
    print(f"Εντολή: {' '.join(command)}")
    try:
        # Χρήση subprocess απευθείας για εμφάνιση εξόδου
        subprocess.run(command, check=True)
        print("Η μετατροπή ffmpeg ολοκληρώθηκε με επιτυχία.")
    except Exception as e:
        print(f"ΣΦΑΛΜΑ ffmpeg: {e}")
        raise Exception(f"Αποτυχία μετατροπής ffmpeg. {e}")
    finally:
        print("Πατήστε Enter για επιστροφή στην εφαρμογή...")
        sys.stdin.read(1)
        stdscr.refresh()

def handle_extraction(stdscr, in_path, out_folder_path, in_ext):
    base_name = os.path.splitext(os.path.basename(in_path))[0]
    extract_path = os.path.join(out_folder_path, base_name)
    os.makedirs(extract_path, exist_ok=True)
    
    if in_ext == '.zip':
        with zipfile.ZipFile(in_path, 'r') as zf:
            zf.extractall(extract_path)
    elif in_ext in ['.tar', '.gz', '.bz2']:
        if in_ext == '.gz' and not in_path.endswith('.tar.gz'):
             out_filename = os.path.splitext(os.path.basename(in_path))[0]
             out_path = os.path.join(out_folder_path, out_filename)
             with gzip.open(in_path, 'rb') as f_in:
                 with open(out_path, 'wb') as f_out:
                     shutil.copyfileobj(f_in, f_out)
             return f"Αποσυμπιέστηκε σε: {out_path}"
        else:
            with tarfile.open(in_path, 'r:*') as tf:
                tf.extractall(extract_path)
    elif in_ext == '.rar':
        if not HAS_UNRAR:
            raise Exception("Δεν βρέθηκε το δυαδικό 'unrar'.")
        import rarfile
        with rarfile.RarFile(in_path) as rf:
            rf.extractall(extract_path)
    elif in_ext == '.7z':
        import py7zr
        with py7zr.SevenZipFile(in_path, 'r') as zf:
            zf.extractall(extract_path)
            
    return f"Εξήχθη σε: {extract_path}"

def handle_data_conversion(stdscr, in_path, out_path, in_ext, out_ext):
    if in_ext == '.csv' and out_ext == '.json':
        data = []
        with open(in_path, 'r', encoding='utf-8', newline='') as f:
            reader = csv.DictReader(f)
            for row in reader:
                data.append(row)
        with open(out_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=4)
    elif in_ext == '.json' and out_ext == '.csv':
        with open(in_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        if not isinstance(data, list) or not data:
            raise Exception("Το JSON πρέπει να είναι μια μη κενή λίστα αντικειμένων.")
        headers = data[0].keys()
        with open(out_path, 'w', encoding='utf-8', newline='') as f:
            writer = csv.DictWriter(f, fieldnames=headers)
            writer.writeheader()
            writer.writerows(data)
    else:
        raise Exception(f"Μετατροπή δεδομένων {in_ext} σε {out_ext} δεν υποστηρίζεται.")

def handle_md_to_html(stdscr, in_path, out_path):
    import markdown
    with open(in_path, 'r', encoding='utf-8') as f:
        html = markdown.markdown(f.read())
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(html)

def handle_single_image_to_pdf(stdscr, in_path, out_path):
    from PIL import Image
    try:
        with Image.open(in_path) as img:
            img_rgb = img.convert('RGB')
            img_rgb.save(out_path, "PDF", resolution=100.0)
    except Exception as e:
        raise Exception(f"Σφάλμα Pillow (Εικόνα->PDF): {e}")

def handle_multi_image_to_pdf(stdscr, image_paths, out_path):
    from PIL import Image
    try:
        images_rgb = []
        for path in image_paths:
            img = Image.open(path)
            images_rgb.append(img.convert('RGB'))
        if not images_rgb:
            raise Exception("Δεν βρέθηκαν εικόνες για μετατροπή.")
        images_rgb[0].save(
            out_path, "PDF", resolution=100.0,
            save_all=True, append_images=images_rgb[1:]
        )
    except Exception as e:
        raise Exception(f"Σφάλμα Pillow (Πολλαπλές Εικόνες->PDF): {e}")

# --- 5. ΚΥΡΙΟΣ ΔΙΑΚΟΜΙΣΤΗΣ ΜΕΤΑΤΡΟΠΩΝ ---

def convert_file(stdscr, in_path, out_folder_name):
    in_ext = os.path.splitext(in_path)[1].lower()
    out_ext = f".{out_folder_name.lower()}"
    base_name = os.path.splitext(os.path.basename(in_path))[0]
    out_folder_path = os.path.join(BASE_CONVERTER_PATH, out_folder_name)
    out_path = os.path.join(out_folder_path, f"{base_name}{out_ext}")

    try:
        if in_ext in ARCHIVE_EXTS:
            out_folder = out_folder_path if in_ext == '.gz' else os.path.join(BASE_CONVERTER_PATH, out_folder_name)
            message = handle_extraction(stdscr, in_path, out_folder, in_ext)
            return (True, message)

        if in_ext in VECTOR_IMAGE_EXTS and out_ext in ['.png', '.pdf']:
            handle_svg_conversion(stdscr, in_path, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")

        if in_ext in LAYERED_IMAGE_EXTS and out_ext in IMAGE_EXTS:
            handle_psd_conversion(stdscr, in_path, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")

        if in_ext in IMAGE_EXTS and out_ext in IMAGE_EXTS:
            handle_image_conversion(stdscr, in_path, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")
            
        if in_ext in IMAGE_EXTS and out_ext == '.pdf':
            handle_single_image_to_pdf(stdscr, in_path, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")

        if in_ext in AV_EXTS and out_ext in AV_EXTS:
            handle_av_conversion(stdscr, in_path, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")
            
        if in_ext in ['.csv', '.json'] and out_ext in ['.csv', '.json']:
            handle_data_conversion(stdscr, in_path, out_path, in_ext, out_ext)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")

        if in_ext == '.md' and out_ext == '.html':
            handle_md_to_html(stdscr, in_path, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")

        if out_ext == '.txt' and in_ext in TEXT_DOC_EXTS:
            text_lines = get_text_from_file(in_path, in_ext)
            with open(out_path, 'w', encoding='utf-8') as f:
                f.writelines(text_lines)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")
            
        if out_ext == '.pdf' and in_ext in TEXT_DOC_EXTS:
            text_lines = get_text_from_file(in_path, in_ext)
            write_text_to_pdf(text_lines, out_path)
            return (True, f"Αποθηκεύτηκε σε: {out_path}")

        return (False, f"Μη υποστηριζόμενη μετατροπή: {in_ext} σε {out_ext}")

    except Exception as e:
        return (False, f"ΣΦΑΛΜΑ: {str(e)}")


# --- 6. ΒΟΗΘΗΤΙΚΕΣ ΣΥΝΑΡΤΗΣΕΙΣ CURSES UI ---

def draw_header(stdscr, title):
    h, w = stdscr.getmaxyx()
    stdscr.attron(curses.color_pair(2))
    stdscr.addstr(0, 0, " " * w, curses.color_pair(2))
    header_text = f" Termux Μετατροπέας Αρχείων (q για έξοδο/πίσω) "
    stdscr.addstr(0, (w - len(header_text)) // 2, header_text, curses.A_REVERSE)
    stdscr.attroff(curses.color_pair(2))
    stdscr.attron(curses.color_pair(1))
    stdscr.addstr(2, (w - len(title)) // 2, title)
    stdscr.attroff(curses.color_pair(1))

def draw_status(stdscr, message, is_error=False):
    h, w = stdscr.getmaxyx()
    nav_hint = " (Βέλη/Enter για Επιλογή, q για Πίσω) "
    hint_len = len(nav_hint)
    color = curses.color_pair(3) if is_error else curses.color_pair(2)
    stdscr.attron(color)
    stdscr.addstr(h - 1, 0, " " * (w - 1))
    stdscr.addstr(h - 1, 1, message[:w - hint_len - 2])
    stdscr.attron(curses.A_REVERSE)
    stdscr.addstr(h - 1, w - hint_len - 1, nav_hint)
    stdscr.attroff(curses.A_REVERSE)
    stdscr.attroff(color)

def run_menu(stdscr, title, options, sub_title=""):
    current_idx = 0
    while True:
        stdscr.clear()
        h, w = stdscr.getmaxyx()
        draw_header(stdscr, title)
        if sub_title:
            stdscr.addstr(4, (w - len(sub_title)) // 2, sub_title)
        draw_status(stdscr, "Επιλέξτε μια επιλογή.")
        
        for i, option in enumerate(options):
            display_option = option
            if len(option) > w - 8:
                display_option = option[:w - 11] + "..."
            x = (w - len(display_option)) // 2 - 2
            y = h // 2 - len(options) // 2 + i
            if i == current_idx:
                stdscr.attron(curses.color_pair(1))
                stdscr.addstr(y, x, f"> {display_option} <")
                stdscr.attroff(curses.color_pair(1))
            else:
                stdscr.addstr(y, x, f"  {display_option}  ")
        
        stdscr.refresh()
        key = stdscr.getch()
        if key == curses.KEY_UP:
            current_idx = (current_idx - 1) % len(options)
        elif key == curses.KEY_DOWN:
            current_idx = (current_idx + 1) % len(options)
        elif key in [curses.KEY_ENTER, 10, 13]:
            return options[current_idx]
        elif key == ord('q'):
            return None

def run_file_selector(stdscr, folder_path, title, input_folder_name):
    try:
        files = [f for f in os.listdir(folder_path) if os.path.isfile(os.path.join(folder_path, f))]
        files.sort()
    except Exception as e:
        draw_status(stdscr, f"Σφάλμα ανάγνωσης {folder_path}: {e}", is_error=True)
        stdscr.getch()
        return None
    
    if not files:
        draw_status(stdscr, f"Δεν βρέθηκαν αρχεία στον {os.path.basename(folder_path)}", is_error=True)
        stdscr.getch()
        return None
        
    options = ["[ .. Επιστροφή .. ]"]
    if input_folder_name in IMAGE_FOLDERS:
        options.append(f"[ ** Μετατροπή ΟΛΩΝ {len(files)} Εικόνων στον '{input_folder_name}' σε ένα PDF ** ]")
    options.extend(files)
    selection = run_menu(stdscr, title, options, f"Φάκελος: /Download/File Converter/{input_folder_name}")
    if selection == "[ .. Επιστροφή .. ]":
        return None
    return selection

def run_confirmation(stdscr, prompt):
    options = ["Ναι", "Όχι"]
    selection = run_menu(stdscr, prompt, options, "Παρακαλώ επιβεβαιώστε")
    return selection

def run_help(stdscr):
    stdscr.clear()
    h, w = stdscr.getmaxyx()
    draw_header(stdscr, "Πώς να Χρησιμοποιήσετε")
    help_text = [
        "1. ΜΕΤΑΚΙΝΗΣΗ ΑΡΧΕΙΩΝ:",
        f"   Μεταβείτε σε: /Download/File Converter/",
        "   Μετακινήστε αρχεία στον σωστό φάκελο ΕΙΣΟΔΟΥ.",
        "",
        "2. ΕΚΤΕΛΕΣΗ ΜΕΤΑΤΡΟΠΕΑ:",
        "   Επιλέξτε 'Μετατροπή Αρχείου' -> Φάκελος ΕΙΣΟΔΟΥ -> Αρχείο.",
        "",
        "3. ΕΠΙΛΟΓΗ ΕΞΟΔΟΥ:",
        "   Επιλέξτε τη μορφή που θέλετε.",
        "",
        "** ΣΗΜΕΙΩΣΕΙΣ **",
        " - Τα αρχεία αρχειοθέτησης εξάγονται στον δικό τους φάκελο.",
        " - Η μετατροπή ήχου/βίντεο χρησιμοποιεί FFmpeg (Αυτόματη εγκατάσταση).",
        " - Οι εικόνες μπορούν να συνδυαστούν σε ένα PDF."
    ]
    for i, line in enumerate(help_text):
        if 5 + i >= h - 2: break
        stdscr.addstr(5 + i, (w - len(line)) // 2, line)
    draw_status(stdscr, "Πατήστε 'q' ή Enter για επιστροφή.")
    stdscr.refresh()
    while True:
        key = stdscr.getch()
        if key in [ord('q'), curses.KEY_ENTER, 10, 13]:
            return

def run_text_input(stdscr, prompt):
    stdscr.clear()
    h, w = stdscr.getmaxyx()
    draw_header(stdscr, "Απαιτείται Εισαγωγή")
    stdscr.addstr(h // 2 - 1, (w - len(prompt)) // 2, prompt)
    box_y, box_x = h // 2 + 1, w // 2 - 20
    stdscr.attron(curses.color_pair(2))
    stdscr.addstr(box_y, box_x, " " * 40)
    stdscr.attroff(curses.color_pair(2))
    draw_status(stdscr, "Πληκτρολογήστε όνομα αρχείου. Enter για Επιβεβαίωση.")
    curses.curs_set(1)
    stdscr.keypad(True)
    text = ""
    while True:
        stdscr.attron(curses.color_pair(2))
        stdscr.addstr(box_y, box_x, " " * 40)
        stdscr.addstr(box_y, box_x + 1, text[:38])
        stdscr.attroff(curses.color_pair(2))
        stdscr.move(box_y, box_x + 1 + len(text))
        stdscr.refresh()
        key = stdscr.getch()
        if key in [curses.KEY_ENTER, 10, 13]:
            break
        elif key == ord('q'):
            text = None; break
        elif key in [curses.KEY_BACKSPACE, 127, 8]:
            text = text[:-1]
        elif 32 <= key <= 126:
            if len(text) < 38: text += chr(key)
    curses.curs_set(0); stdscr.keypad(False); return text

def run_multi_image_to_pdf_wizard(stdscr, input_folder_path, input_folder_name):
    try:
        image_paths = [
            os.path.join(input_folder_path, f) 
            for f in os.listdir(input_folder_path) 
            if os.path.splitext(f)[1].lower() in IMAGE_EXTS
        ]
        image_paths.sort()
    except Exception as e:
        draw_status(stdscr, f"Σφάλμα: {e}", is_error=True); stdscr.getch(); return
    if not image_paths:
        draw_status(stdscr, "Δεν βρέθηκαν εικόνες.", is_error=True); stdscr.getch(); return
    confirm = run_confirmation(stdscr, f"Συνδυασμός {len(image_paths)} εικόνων σε ένα PDF;")
    if confirm != "Ναι": return
    default_name = f"{input_folder_name}_Album"
    filename = run_text_input(stdscr, f"Εισάγετε όνομα PDF (προεπιλογή: {default_name})")
    if filename is None: return
    if not filename: filename = default_name
    out_folder_path = os.path.join(BASE_CONVERTER_PATH, "PDF")
    out_path = os.path.join(out_folder_path, f"{filename}.pdf")
    draw_status(stdscr, "Εργασία..."); stdscr.refresh()
    try:
        handle_multi_image_to_pdf(stdscr, image_paths, out_path)
        draw_status(stdscr, f"Επιτυχία! Αποθηκεύτηκε σε: /PDF/{filename}.pdf")
    except Exception as e:
        draw_status(stdscr, f"ΣΦΑΛΜΑ: {e}", is_error=True)
    stdscr.getch()

# --- 7. ΚΥΡΙΑ ΕΦΑΡΜΟΓΗ CURSES ---

def main(stdscr):
    curses.curs_set(0); stdscr.nodelay(0); stdscr.timeout(-1)
    curses.start_color()
    curses.init_pair(1, curses.COLOR_CYAN, curses.COLOR_BLACK)
    curses.init_pair(2, curses.COLOR_BLACK, curses.COLOR_WHITE)
    curses.init_pair(3, curses.COLOR_WHITE, curses.COLOR_RED)

    while True:
        main_choice = run_menu(stdscr, "Κύριο Μενού", ["Μετατροπή Αρχείου", "Βοήθεια / Πώς να Χρησιμοποιήσετε", "Έξοδος"])
        if main_choice == "Έξοδος" or main_choice is None: break
        if main_choice == "Βοήθεια / Πώς να Χρησιμοποιήσετε": run_help(stdscr); continue
            
        input_folder = run_menu(stdscr, "Βήμα 1: Επιλογή Φακέλου ΕΙΣΟΔΟΥ", FOLDER_NAMES)
        if input_folder is None: continue
        input_folder_path = os.path.join(BASE_CONVERTER_PATH, input_folder)
        input_file = run_file_selector(stdscr, input_folder_path, f"Βήμα 2: Επιλογή Αρχείου από '{input_folder}'", input_folder)
        if input_file is None: continue
        
        if input_file.startswith("[ ** Μετατροπή ΟΛΩΝ"):
            run_multi_image_to_pdf_wizard(stdscr, input_folder_path, input_folder)
            continue
            
        full_input_path = os.path.join(input_folder_path, input_file)
        in_ext = os.path.splitext(input_file)[1].lower()
        if in_ext in ARCHIVE_EXTS:
            output_folder = input_folder
            prompt = f"Εξαγωγή '{input_file}' εδώ;"
        else:
            output_folder = run_menu(stdscr, "Βήμα 3: Επιλογή Μορφής ΕΞΟΔΟΥ", FOLDER_NAMES)
            if output_folder is None: continue
            if output_folder == input_folder:
                draw_status(stdscr, "Η είσοδος και η έξοδος δεν μπορούν να είναι ίδιες.", is_error=True)
                stdscr.getch(); continue
            prompt = f"Μετατροπή '{input_file}' σε {output_folder};"
             
        confirm = run_confirmation(stdscr, prompt)
        if confirm != "Ναι": continue

        draw_status(stdscr, "Εργασία, παρακαλώ περιμένετε..."); stdscr.refresh()
        success, message = convert_file(stdscr, full_input_path, output_folder)
        draw_status(stdscr, message, is_error=not success)
        stdscr.getch()

# --- 8. ΣΗΜΕΙΟ ΕΝΑΡΞΗΣ SCRIPT ---

if __name__ == "__main__":
    try:
        clear_screen_standard()
        print("--- Αρχικοποίηση Termux Μετατροπέα v3.2 (Αυτόματη Εγκατάσταση) ---")
        
        # 1. Εγκατάσταση Δυαδικών Συστήματος (pkg)
        install_termux_system_deps()
        
        # 2. Εγκατάσταση Modules Python (pip)
        check_and_install_python_deps()
        
        # 3. Επαλήθευση Ολοκληρωμάτων
        check_external_bins_status()
        
        # 4. Έλεγχος Αποθήκευσης
        ensure_storage_access()
        
        setup_folders()
        
        print("--- Η Διαμόρφωση Ολοκληρώθηκε ---")
        print(f"Οι φάκελοι είναι έτοιμοι σε: /storage/emulated/0/Download/File Converter/")
        print("\nΕκκίνηση εφαρμογής...")
        time.sleep(1)
        
        curses.wrapper(main)
        print("Ο Μετατροπέας Αρχείων τερμάτισε επιτυχώς.")

    except KeyboardInterrupt:
        print("\nΈξοδος...")
    except Exception as e:
        print("\nΠροέκυψε κρίσιμο σφάλμα:")
        traceback.print_exc()
    finally:
        os.system('clear')