#!/usr/bin/env python

"""
An advanced, interactive file converter for Termux using 'curses'.
Version 3.1: Organizes all folders into a master "File Converter" directory.

This script will:
1. Check and install required Python libraries.
2. Check for external binaries (ffmpeg, unrar, cairo).
3. Create a "File Converter" folder in Downloads, containing 40 organizer sub-folders.
4. Provide a 'curses'-based UI for navigation.
5. Support many document, image, A/V, and data conversions.

CRITICAL SETUP (in Termux):
pkg install ffmpeg unrar libcairo libgirepository libjpeg-turbo libpng
"""

import sys
import os
import subprocess
import importlib.util
import time
import curses
import traceback
import zipfile
import tarfile
import csv
import json
import gzip
import shutil
from contextlib import redirect_stderr, redirect_stdout

# --- 1. SETUP & CONFIGURATION ---

# (14) Python libraries to auto-install
REQUIRED_MODULES = {
    "Pillow": "Pillow",         # Images
    "reportlab": "reportlab",   # PDF creation
    "docx": "python-docx",    # Word docs
    "odf": "odfpy",           # OpenOffice docs
    "bs4": "beautifulsoup4",  # HTML/XML parsing
    "markdown": "Markdown",     # Markdown parsing
    "lxml": "lxml",           # XML/HTML parsing
    "cairosvg": "cairosvg",     # SVG conversion
    "psd_tools": "psd-tools",   # PSD reading
    "striprtf": "striprtf",     # RTF reading
    "EbookLib": "EbookLib",     # EPUB reading
    "pptx": "python-pptx",    # PowerPoint reading
    "rarfile": "rarfile",       # RAR extraction
    "py7zr": "py7zr"          # 7-Zip extraction
}

# (40) Folders to create
FOLDER_NAMES = [
    # Images (10)
    "JPG", "PNG", "WEBP", "BMP", "TIFF", "GIF", "ICO", "TGA", "SVG", "PSD",
    # Documents (12)
    "PDF", "TXT", "DOCX", "ODT", "HTML", "MD", "CSV", "RTF", "EPUB", "JSON", "XML", "PPTX",
    # Archives (5)
    "ZIP", "TAR", "RAR", "7Z", "GZ",
    # Audio (7)
    "MP3", "WAV", "OGG", "FLAC", "M4A", "AAC", "WMA",
    # Video (6)
    "MP4", "MKV", "AVI", "MOV", "WMV", "FLV"
]

# Helper lists for logic
IMAGE_FOLDERS = ["JPG", "PNG", "WEBP", "BMP", "TIFF", "GIF", "ICO", "TGA", "SVG", "PSD"]
IMAGE_EXTS = ['.jpg', '.jpeg', '.png', '.webp', '.bmp', '.tiff', '.gif', '.ico', '.tga']
VECTOR_IMAGE_EXTS = ['.svg']
LAYERED_IMAGE_EXTS = ['.psd']
AV_EXTS = ['.mp3', '.wav', '.ogg', '.flac', '.m4a', '.aac', '.wma', '.mp4', '.mkv', '.avi', '.mov', '.wmv', '.flv']
ARCHIVE_EXTS = ['.zip', '.tar', '.gz', '.bz2', '.rar', '.7z']
TEXT_DOC_EXTS = ['.txt', '.docx', '.odt', '.html', '.md', '.csv', '.rtf', '.epub', '.json', '.xml', '.pptx', '.svg']
DATA_EXTS = ['.csv', '.json', '.xml']

# Paths
STORAGE_PATH = "/storage/emulated/0"
DOWNLOAD_PATH = os.path.join(STORAGE_PATH, "Download")
# --- MODIFIED: All folders are now nested inside "File Converter" ---
BASE_CONVERTER_PATH = os.path.join(DOWNLOAD_PATH, "File Converter")

# Global flags for external binaries
HAS_FFMPEG = False
HAS_UNRAR = False
HAS_CAIRO = False

# --- 2. PRE-CURSES SETUP FUNCTIONS (Standard Print) ---

def clear_screen_standard():
    os.system('clear')

def check_and_install_dependencies():
    """Checks and installs required Python modules."""
    print("--- Checking Required Python Libraries (14) ---")
    all_installed = True
    for module_name, package_name in REQUIRED_MODULES.items():
        spec = importlib.util.find_spec(module_name)
        if spec is None:
            all_installed = False
            print(f"Installing '{package_name}'...")
            try:
                with open(os.devnull, 'w') as devnull:
                    with redirect_stdout(devnull), redirect_stderr(devnull):
                        subprocess.run([sys.executable, "-m", "pip", "install", package_name], check=True)
                print(f"Successfully installed '{package_name}'.")
            except Exception:
                print(f"ERROR: Failed to install '{package_name}'.")
                print(f"Please install it manually: pip install {package_name}")
                sys.exit(1)
        else:
            # print(f"Library '{package_name}' is already installed.")
            pass
    
    if all_installed:
        print("All Python libraries are present.\n")
    else:
        print("All required libraries are now installed.\n")
    time.sleep(0.5)

def check_external_bins():
    """Checks for 'ffmpeg', 'unrar', and 'cairo'."""
    global HAS_FFMPEG, HAS_UNRAR, HAS_CAIRO
    print("--- Checking External Binaries ---")
    
    # Check ffmpeg
    try:
        with open(os.devnull, 'w') as devnull:
            subprocess.run(["ffmpeg", "-version"], check=True, stdout=devnull, stderr=devnull)
        print("Found 'ffmpeg'. Audio/Video conversions are ENABLED.")
        HAS_FFMPEG = True
    except (subprocess.CalledProcessError, FileNotFoundError):
        print("WARNING: 'ffmpeg' not found. A/V conversions DISABLED.")
        print("  To enable, run: pkg install ffmpeg\n")
        HAS_FFMPEG = False

    # Check unrar
    try:
        with open(os.devnull, 'w') as devnull:
            subprocess.run(["unrar"], check=True, stdout=devnull, stderr=devnull)
        print("Found 'unrar'. RAR extraction is ENABLED.")
        HAS_UNRAR = True
    except (subprocess.CalledProcessError, FileNotFoundError):
        print("WARNING: 'unrar' not found. RAR extraction DISABLED.")
        print("  To enable, run: pkg install unrar\n")
        HAS_UNRAR = False
        
    # Check cairo (for SVG)
    if importlib.util.find_spec("cairosvg") is not None:
        print("Found 'cairosvg'. SVG conversions are ENABLED.")
        HAS_CAIRO = True
    else:
        print("WARNING: 'cairosvg' Python lib not found. SVG conversion DISABLED.")
        print("  The script tried to install it, but it may have failed.")
        print("  You may also need: pkg install libcairo libgirepository\n")
        HAS_CAIRO = False
        
    print("")
    time.sleep(0.5)


def check_storage_access():
    print("--- Checking Storage Access ---")
    if not os.path.exists(DOWNLOAD_PATH):
        print(f"ERROR: Cannot access internal storage at '{DOWNLOAD_PATH}'.")
        print("Please run 'termux-setup-storage' in your Termux terminal,")
        print("grant the permission, and then run this script again.")
        sys.exit(1)
    print("Storage access confirmed.\n")
    time.sleep(0.5)

def setup_folders():
    # --- MODIFIED: Creates the main "File Converter" directory first ---
    print(f"--- Setting up Organizer Folders ---")
    print(f"Location: {BASE_CONVERTER_PATH}")
    try:
        # 1. Create the main parent directory
        os.makedirs(BASE_CONVERTER_PATH, exist_ok=True)
        # 2. Create all 40 sub-folders inside it
        for folder in FOLDER_NAMES:
            os.makedirs(os.path.join(BASE_CONVERTER_PATH, folder), exist_ok=True)
        print(f"Successfully created/verified {len(FOLDER_NAMES)} sub-folders.\n")
    except Exception as e:
        print(f"ERROR: Could not create folders: {e}")
        sys.exit(1)
    time.sleep(0.5)

# --- 3. IMPORTS (Post-Installation) ---
try:
    from PIL import Image
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from docx import Document
    from odf.opendocument import load as odf_load
    from odf.text import P as odf_P
    from bs4 import BeautifulSoup
    import markdown
    import lxml
    import cairosvg
    from psd_tools import PSDImage
    from striprtf.striprtf import rtf_to_text
    from ebooklib import epub, ITEM_DOCUMENT
    import pptx
    import rarfile
    import py7zr
except ImportError as e:
    print(f"CRITICAL ERROR: Failed to import a library: {e}")
    print("Please ensure all dependencies are installed (see startup logs).")
    sys.exit(1)

# --- 4. CORE CONVERSION LOGIC ---

def get_text_from_file(input_path, in_ext):
    """Helper to extract raw text from various document types."""
    text_lines = []
    try:
        if in_ext == '.txt':
            with open(input_path, 'r', encoding='utf-8') as f:
                text_lines = f.readlines()
        elif in_ext == '.docx':
            doc = Document(input_path)
            text_lines = [para.text + '\n' for para in doc.paragraphs]
        elif in_ext == '.odt':
            doc = odf_load(input_path)
            for para in doc.getElementsByType(odf_P):
                text_lines.append(str(para) + '\n')
        elif in_ext in ['.html', '.xml', '.svg']:
            with open(input_path, 'r', encoding='utf-8') as f:
                parser = 'lxml' if in_ext != '.html' else 'html.parser'
                soup = BeautifulSoup(f, parser)
                text_lines = [line + '\n' for line in soup.get_text().splitlines()]
        elif in_ext == '.md':
            with open(input_path, 'r', encoding='utf-8') as f:
                html = markdown.markdown(f.read())
                soup = BeautifulSoup(html, 'html.parser')
                text_lines = [line + '\n' for line in soup.get_text().splitlines()]
        elif in_ext == '.csv':
            with open(input_path, 'r', encoding='utf-8', newline='') as f:
                reader = csv.reader(f)
                text_lines = [','.join(row) + '\n' for row in reader]
        elif in_ext == '.json':
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                text_lines = [json.dumps(data, indent=2)]
        elif in_ext == '.rtf':
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            text_lines = [rtf_to_text(content)]
        elif in_ext == '.epub':
            book = epub.read_epub(input_path)
            for item in book.get_items():
                if item.get_type() == ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text_lines.append(soup.get_text() + '\n\n') # Add space between chapters
        elif in_ext == '.pptx':
            prs = pptx.Presentation(input_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_lines.append(shape.text + '\n')
    except Exception as e:
        raise Exception(f"Text extraction failed: {e}")
    return text_lines

def write_text_to_pdf(text_lines, output_path):
    # (Unchanged)
    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    margin_x, margin_y = 0.75 * inch, 1 * inch
    text_object = c.beginText(margin_x, height - margin_y)
    text_object.setFont("Helvetica", 10)
    line_height, y = 12, height - margin_y
    for line in text_lines:
        for sub_line in line.split('\n'): # Handle multi-line strings
            if y < margin_y:
                c.drawText(text_object)
                c.showPage()
                text_object = c.beginText(margin_x, height - margin_y)
                text_object.setFont("Helvetica", 10)
                y = height - margin_y
            text_object.textLine(sub_line.strip('\r'))
            y -= line_height
    c.drawText(text_object)
    c.save()

def handle_image_conversion(stdscr, in_path, out_path):
    # (Pillow handler, unchanged)
    with Image.open(in_path) as img:
        if out_path.lower().endswith(('.jpg', '.jpeg')):
            if img.mode == 'RGBA':
                img = img.convert('RGB')
        img.save(out_path)

def handle_svg_conversion(stdscr, in_path, out_path):
    """Converts SVG to PNG or PDF."""
    if not HAS_CAIRO:
        raise Exception("Cairo/SVG libraries not installed.")
    out_ext = os.path.splitext(out_path)[1].lower()
    if out_ext == '.png':
        cairosvg.svg2png(url=in_path, write_to=out_path)
    elif out_ext == '.pdf':
        cairosvg.svg2pdf(url=in_path, write_to=out_path)
    else:
        raise Exception(f"SVG conversion to {out_ext} not supported.")

def handle_psd_conversion(stdscr, in_path, out_path):
    """Converts PSD composite to a flat image."""
    psd = PSDImage.open(in_path)
    composite_image = psd.composite()
    composite_image.save(out_path)

def handle_av_conversion(stdscr, in_path, out_path):
    # (Unchanged)
    if not HAS_FFMPEG:
        raise Exception("'ffmpeg' not found. A/V conversion is disabled.")
    command = ['ffmpeg', '-i', in_path, '-y', out_path]
    curses.endwin()
    print("--- Running ffmpeg ---")
    print(f"Command: {' '.join(command)}")
    print("This may take some time...")
    try:
        with open(os.devnull, 'w') as devnull:
            subprocess.run(command, check=True, stdout=devnull, stderr=subprocess.STDOUT)
        print("ffmpeg conversion finished successfully.")
    except Exception as e:
        print(f"ffmpeg ERROR: {e}")
        raise Exception(f"ffmpeg conversion failed. {e}")
    finally:
        print("Press Enter to return to the app...")
        sys.stdin.read(1)
        stdscr.refresh()

def handle_extraction(stdscr, in_path, out_folder_path, in_ext):
    """Extracts various archive types."""
    base_name = os.path.splitext(os.path.basename(in_path))[0]
    extract_path = os.path.join(out_folder_path, base_name)
    os.makedirs(extract_path, exist_ok=True)
    
    if in_ext == '.zip':
        with zipfile.ZipFile(in_path, 'r') as zf:
            zf.extractall(extract_path)
    elif in_ext in ['.tar', '.gz', '.bz2']:
        if in_ext == '.gz' and not in_path.endswith('.tar.gz'): # Single file gzip
             out_filename = os.path.splitext(os.path.basename(in_path))[0]
             out_path = os.path.join(out_folder_path, out_filename) # Extract to folder, not subfolder
             with gzip.open(in_path, 'rb') as f_in:
                 with open(out_path, 'wb') as f_out:
                     shutil.copyfileobj(f_in, f_out)
             return f"Decompressed to: {out_path}" # Return different message
        else: # .tar, .tar.gz, .tar.bz2
            with tarfile.open(in_path, 'r:*') as tf:
                tf.extractall(extract_path)
    elif in_ext == '.rar':
        if not HAS_UNRAR:
            raise Exception("'unrar' binary not found.")
        with rarfile.RarFile(in_path) as rf:
            rf.extractall(extract_path)
    elif in_ext == '.7z':
        with py7zr.SevenZipFile(in_path, 'r') as zf:
            zf.extractall(extract_path)
            
    return f"Extracted to: {extract_path}" # Default success message

def handle_data_conversion(stdscr, in_path, out_path, in_ext, out_ext):
    """Handles CSV <-> JSON conversions."""
    if in_ext == '.csv' and out_ext == '.json':
        data = []
        with open(in_path, 'r', encoding='utf-8', newline='') as f:
            reader = csv.DictReader(f)
            for row in reader:
                data.append(row)
        with open(out_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=4)
    elif in_ext == '.json' and out_ext == '.csv':
        with open(in_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        if not isinstance(data, list) or not data:
            raise Exception("JSON must be a non-empty list of objects.")
        if not all(isinstance(x, dict) for x in data):
            raise Exception("JSON must be a list of objects (dicts).")
            
        headers = data[0].keys()
        with open(out_path, 'w', encoding='utf-8', newline='') as f:
            writer = csv.DictWriter(f, fieldnames=headers)
            writer.writeheader()
            writer.writerows(data)
    else:
        raise Exception(f"Data conversion {in_ext} to {out_ext} not supported.")

def handle_md_to_html(stdscr, in_path, out_path):
    # (Unchanged)
    with open(in_path, 'r', encoding='utf-8') as f:
        html = markdown.markdown(f.read())
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(html)

def handle_single_image_to_pdf(stdscr, in_path, out_path):
    # (Unchanged)
    try:
        with Image.open(in_path) as img:
            img_rgb = img.convert('RGB')
            img_rgb.save(out_path, "PDF", resolution=100.0)
    except Exception as e:
        raise Exception(f"Pillow (Image->PDF) Error: {e}")

def handle_multi_image_to_pdf(stdscr, image_paths, out_path):
    # (Unchanged)
    try:
        images_rgb = []
        for path in image_paths:
            img = Image.open(path)
            images_rgb.append(img.convert('RGB'))
        if not images_rgb:
            raise Exception("No images found to convert.")
        images_rgb[0].save(
            out_path, "PDF", resolution=100.0,
            save_all=True, append_images=images_rgb[1:]
        )
    except Exception as e:
        raise Exception(f"Pillow (Multi-Image->PDF) Error: {e}")

# --- 5. MAIN CONVERSION ROUTER ---

def convert_file(stdscr, in_path, out_folder_name):
    """
    Main router function for dispatching conversion tasks.
    Returns (success_bool, message_string)
    """
    in_ext = os.path.splitext(in_path)[1].lower()
    out_ext = f".{out_folder_name.lower()}"
    
    base_name = os.path.splitext(os.path.basename(in_path))[0]
    out_folder_path = os.path.join(BASE_CONVERTER_PATH, out_folder_name)
    out_path = os.path.join(out_folder_path, f"{base_name}{out_ext}")

    try:
        # --- Route 1: Extraction ---
        if in_ext in ARCHIVE_EXTS:
            # Note: GZ files will be decompressed *into* the folder named GZ
            # All others (ZIP, TAR, RAR, 7Z) extract to a *sub-folder*
            out_folder = out_folder_path if in_ext == '.gz' else os.path.join(BASE_CONVERTER_PATH, out_folder_name)
            message = handle_extraction(stdscr, in_path, out_folder, in_ext)
            return (True, message)

        # --- Route 2: SVG Conversion (to PNG, PDF) ---
        if in_ext in VECTOR_IMAGE_EXTS and out_ext in ['.png', '.pdf']:
            handle_svg_conversion(stdscr, in_path, out_path)
            return (True, f"Saved to: {out_path}")

        # --- Route 3: PSD Conversion (to flat image) ---
        if in_ext in LAYERED_IMAGE_EXTS and out_ext in IMAGE_EXTS:
            handle_psd_conversion(stdscr, in_path, out_path)
            return (True, f"Saved to: {out_path}")

        # --- Route 4: Image-to-Image (Pillow) ---
        if in_ext in IMAGE_EXTS and out_ext in IMAGE_EXTS:
            handle_image_conversion(stdscr, in_path, out_path)
            return (True, f"Saved to: {out_path}")
            
        # --- Route 5: Single Image-to-PDF ---
        if in_ext in IMAGE_EXTS and out_ext == '.pdf':
            handle_single_image_to_pdf(stdscr, in_path, out_path)
            return (True, f"Saved to: {out_path}")

        # --- Route 6: A/V-to-A/V (ffmpeg) ---
        if in_ext in AV_EXTS and out_ext in AV_EXTS:
            handle_av_conversion(stdscr, in_path, out_path)
            return (True, f"Saved to: {out_path}")
            
        # --- Route 7: Data Conversion (CSV <-> JSON) ---
        if in_ext in ['.csv', '.json'] and out_ext in ['.csv', '.json']:
            handle_data_conversion(stdscr, in_path, out_path, in_ext, out_ext)
            return (True, f"Saved to: {out_path}")

        # --- Route 8: MD-to-HTML ---
        if in_ext == '.md' and out_ext == '.html':
            handle_md_to_html(stdscr, in_path, out_path)
            return (True, f"Saved to: {out_path}")

        # --- Route 9: Anything-to-TXT ---
        if out_ext == '.txt' and in_ext in TEXT_DOC_EXTS:
            text_lines = get_text_from_file(in_path, in_ext)
            with open(out_path, 'w', encoding='utf-8') as f:
                f.writelines(text_lines)
            return (True, f"Saved to: {out_path}")
            
        # --- Route 10: Anything-to-PDF ---
        if out_ext == '.pdf' and in_ext in TEXT_DOC_EXTS:
            text_lines = get_text_from_file(in_path, in_ext)
            write_text_to_pdf(text_lines, out_path)
            return (True, f"Saved to: {out_path}")

        # --- No Route Found ---
        return (False, f"Unsupported conversion: {in_ext} to {out_ext}")

    except Exception as e:
        return (False, f"ERROR: {str(e)}")


# --- 6. CURSES UI HELPER FUNCTIONS ---
# (Unchanged from v2, but added note to run_help)

def draw_header(stdscr, title):
    h, w = stdscr.getmaxyx()
    stdscr.attron(curses.color_pair(2))
    stdscr.addstr(0, 0, " " * w, curses.color_pair(2))
    header_text = f" Termux Curses Converter (q to quit/back) "
    stdscr.addstr(0, (w - len(header_text)) // 2, header_text, curses.A_REVERSE)
    stdscr.attroff(curses.color_pair(2))
    stdscr.attron(curses.color_pair(1))
    stdscr.addstr(2, (w - len(title)) // 2, title)
    stdscr.attroff(curses.color_pair(1))

def draw_status(stdscr, message, is_error=False):
    h, w = stdscr.getmaxyx()
    nav_hint = " (Use Arrows, Enter to Select, q to Back) "
    hint_len = len(nav_hint)
    color = curses.color_pair(3) if is_error else curses.color_pair(2)
    stdscr.attron(color)
    stdscr.addstr(h - 1, 0, " " * (w - 1))
    stdscr.addstr(h - 1, 1, message[:w - hint_len - 2])
    stdscr.attron(curses.A_REVERSE)
    stdscr.addstr(h - 1, w - hint_len - 1, nav_hint)
    stdscr.attroff(curses.A_REVERSE)
    stdscr.attroff(color)

def run_menu(stdscr, title, options, sub_title=""):
    current_idx = 0
    while True:
        stdscr.clear()
        h, w = stdscr.getmaxyx()
        draw_header(stdscr, title)
        if sub_title:
            stdscr.addstr(4, (w - len(sub_title)) // 2, sub_title)
        draw_status(stdscr, "Select an option.")
        
        for i, option in enumerate(options):
            display_option = option
            if len(option) > w - 8:
                display_option = option[:w - 11] + "..."
            x = (w - len(display_option)) // 2 - 2
            y = h // 2 - len(options) // 2 + i
            if i == current_idx:
                stdscr.attron(curses.color_pair(1))
                stdscr.addstr(y, x, f"> {display_option} <")
                stdscr.attroff(curses.color_pair(1))
            else:
                stdscr.addstr(y, x, f"  {display_option}  ")
        
        stdscr.refresh()
        key = stdscr.getch()
        if key == curses.KEY_UP:
            current_idx = (current_idx - 1) % len(options)
        elif key == curses.KEY_DOWN:
            current_idx = (current_idx + 1) % len(options)
        elif key in [curses.KEY_ENTER, 10, 13]:
            return options[current_idx]
        elif key == ord('q'):
            return None

def run_file_selector(stdscr, folder_path, title, input_folder_name):
    try:
        files = [f for f in os.listdir(folder_path) if os.path.isfile(os.path.join(folder_path, f))]
        files.sort()
    except Exception as e:
        draw_status(stdscr, f"Error reading {folder_path}: {e}", is_error=True)
        stdscr.getch()
        return None
    
    if not files:
        draw_status(stdscr, f"No files found in {os.path.basename(folder_path)}", is_error=True)
        stdscr.getch()
        return None
        
    options = ["[ .. Go Back .. ]"]
    
    if input_folder_name in IMAGE_FOLDERS:
        options.append(f"[ ** Convert ALL {len(files)} Images in '{input_folder_name}' to one PDF ** ]")
    
    options.extend(files)
    # --- MODIFIED: Updated path in subtitle ---
    selection = run_menu(stdscr, title, options, f"Folder: /Download/File Converter/{input_folder_name}")
    if selection == "[ .. Go Back .. ]":
        return None
    return selection

def run_confirmation(stdscr, prompt):
    options = ["Yes", "No"]
    selection = run_menu(stdscr, prompt, options, "Please confirm")
    return selection

def run_help(stdscr):
    stdscr.clear()
    h, w = stdscr.getmaxyx()
    draw_header(stdscr, "How to Use")
    help_text = [
        "This converter uses a simple 3-step process:",
        "",
        "1. MOVE YOUR FILES:",
        "   Use your phone's File Manager. Go to:",
        # --- MODIFIED: Updated path in help text ---
        f"   /Download/File Converter/",
        "   Move files into the correct folder (e.g., move",
        "   'report.docx' into the 'DOCX' folder).",
        "",
        "2. RUN THIS CONVERTER:",
        "   Select 'Convert a File' from the main menu.",
        "",
        "3. FOLLOW THE STEPS:",
        "   Step 1: Choose the INPUT folder (e.g., 'DOCX').",
        "   Step 2: Choose the file you want to convert.",
        "   Step 3: Choose the OUTPUT format (e.g., 'PDF').",
        "",
        "** SPECIAL CONVERSIONS **",
        " - Archives (ZIP, RAR, 7Z, TAR): Choose 'ZIP' -> 'file.zip' -> 'ZIP'",
        "   This will extract 'file.zip' into a new folder: /ZIP/file/",
        " - Multi-Image PDF: Choose 'JPG' -> '[ ** Convert ALL... ** ]'",
        "   This combines all images in 'JPG' into one PDF.",
        " - Data: You can convert CSV <-> JSON.",
        " - A/V: Audio/Video conversions require 'ffmpeg' (see startup)."
    ]
    for i, line in enumerate(help_text):
        if 5 + i >= h - 2: break # Stop if too long for screen
        stdscr.addstr(5 + i, (w - len(line)) // 2, line)
    draw_status(stdscr, "Press 'q' or Enter to go back.")
    stdscr.refresh()
    while True:
        key = stdscr.getch()
        if key in [ord('q'), curses.KEY_ENTER, 10, 13]:
            return

def run_text_input(stdscr, prompt):
    # (Unchanged)
    stdscr.clear()
    h, w = stdscr.getmaxyx()
    draw_header(stdscr, "Input Required")
    stdscr.addstr(h // 2 - 1, (w - len(prompt)) // 2, prompt)
    box_y, box_x = h // 2 + 1, w // 2 - 20
    stdscr.attron(curses.color_pair(2))
    stdscr.addstr(box_y, box_x, " " * 40)
    stdscr.attroff(curses.color_pair(2))
    draw_status(stdscr, "Type the filename (no extension). Press Enter when done.")
    curses.curs_set(1)
    stdscr.keypad(True)
    text = ""
    while True:
        stdscr.attron(curses.color_pair(2))
        stdscr.addstr(box_y, box_x, " " * 40)
        stdscr.addstr(box_y, box_x + 1, text[:38])
        stdscr.attroff(curses.color_pair(2))
        stdscr.move(box_y, box_x + 1 + len(text))
        stdscr.refresh()
        key = stdscr.getch()
        if key in [curses.KEY_ENTER, 10, 13]:
            break
        elif key == ord('q'):
            text = None; break
        elif key in [curses.KEY_BACKSPACE, 127, 8]:
            text = text[:-1]
        elif 32 <= key <= 126:
            if len(text) < 38: text += chr(key)
    curses.curs_set(0); stdscr.keypad(False); return text

def run_multi_image_to_pdf_wizard(stdscr, input_folder_path, input_folder_name):
    # (Unchanged)
    try:
        image_paths = [
            os.path.join(input_folder_path, f) 
            for f in os.listdir(input_folder_path) 
            if os.path.splitext(f)[1].lower() in IMAGE_EXTS
        ]
        image_paths.sort()
    except Exception as e:
        draw_status(stdscr, f"Error reading images: {e}", is_error=True); stdscr.getch(); return
    if not image_paths:
        draw_status(stdscr, "No images found in that folder.", is_error=True); stdscr.getch(); return
    confirm = run_confirmation(stdscr, f"Combine all {len(image_paths)} images in '{input_folder_name}' into one PDF?")
    if confirm != "Yes": return
    default_name = f"{input_folder_name}_Album"
    filename = run_text_input(stdscr, f"Enter a name for the PDF (default: {default_name})")
    if filename is None: return
    if not filename: filename = default_name
    out_folder_path = os.path.join(BASE_CONVERTER_PATH, "PDF")
    out_path = os.path.join(out_folder_path, f"{filename}.pdf")
    draw_status(stdscr, "Working... combining images into PDF..."); stdscr.refresh()
    try:
        handle_multi_image_to_pdf(stdscr, image_paths, out_path)
        draw_status(stdscr, f"Success! Saved to: /PDF/{filename}.pdf")
    except Exception as e:
        draw_status(stdscr, f"ERROR: {e}", is_error=True)
    stdscr.getch()

# --- 7. MAIN CURSES APPLICATION ---

def main(stdscr):
    # (Unchanged)
    curses.curs_set(0); stdscr.nodelay(0); stdscr.timeout(-1)
    curses.start_color()
    curses.init_pair(1, curses.COLOR_CYAN, curses.COLOR_BLACK) # Highlight
    curses.init_pair(2, curses.COLOR_BLACK, curses.COLOR_WHITE) # Header/Footer
    curses.init_pair(3, curses.COLOR_WHITE, curses.COLOR_RED)   # Error

    while True:
        main_choice = run_menu(stdscr, "Main Menu", ["Convert a File", "Help / How to Use", "Exit"])
        if main_choice == "Exit" or main_choice is None: break
        if main_choice == "Help / How to Use": run_help(stdscr); continue
            
        input_folder = run_menu(stdscr, "Step 1: Choose INPUT Folder", FOLDER_NAMES)
        if input_folder is None: continue
        input_folder_path = os.path.join(BASE_CONVERTER_PATH, input_folder)
        input_file = run_file_selector(stdscr, input_folder_path, f"Step 2: Choose File from '{input_folder}'", input_folder)
        if input_file is None: continue
        
        if input_file.startswith("[ ** Convert ALL"):
            run_multi_image_to_pdf_wizard(stdscr, input_folder_path, input_folder)
            continue
            
        full_input_path = os.path.join(input_folder_path, input_file)
        
        # --- Special Case: Archive Extraction ---
        in_ext = os.path.splitext(input_file)[1].lower()
        if in_ext in ARCHIVE_EXTS:
            # For archives, the "output folder" is just the folder of the same type
            # e.g., extract a ZIP file into the "ZIP" folder.
            output_folder = input_folder
            prompt = f"Extract '{input_file}' into '/{output_folder}/'?"
            if in_ext not in ['.zip', '.rar', '.7z', '.tar', '.gz', '.bz2']:
                output_folder = input_folder # Failsafe
                
        else:
            # --- Normal Conversion Path ---
            output_folder = run_menu(stdscr, "Step 3: Choose OUTPUT Format/Folder", FOLDER_NAMES)
            if output_folder is None: continue
            if output_folder == input_folder:
                draw_status(stdscr, "Error: Input and Output folder cannot be the same.", is_error=True)
                stdscr.getch(); continue
            prompt = f"Convert '{input_file}' to {output_folder} format?"
             
        confirm = run_confirmation(stdscr, prompt)
        if confirm != "Yes": continue

        draw_status(stdscr, "Working, please wait..."); stdscr.refresh()
        success, message = convert_file(stdscr, full_input_path, output_folder)
        draw_status(stdscr, message, is_error=not success)
        stdscr.getch()

# --- 8. SCRIPT ENTRYPOINT ---

if __name__ == "__main__":
    try:
        clear_screen_standard()
        print("--- Initializing Termux Converter v3.1 (40 Formats) ---")
        check_and_install_dependencies()
        check_external_bins()
        check_storage_access()
        setup_folders()
        
        print("--- Setup Complete ---")
        # --- MODIFIED: Updated final path message ---
        print(f"Organizer folders are ready in: /storage/emulated/0/Download/File Converter/")
        print("\nStarting application...")
        time.sleep(1)
        
        curses.wrapper(main)
        print("File Converter exited successfully.")

    except KeyboardInterrupt:
        print("\nExiting...")
    except Exception as e:
        print("\nA critical error occurred:")
        traceback.print_exc()
    finally:
        os.system('clear')